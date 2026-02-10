"""Document extraction utilities."""

import io
import os
import warnings
from pathlib import Path

import chardet
import numpy as np
import pypdf
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from odf import teletype
from odf import text as odf_text
from odf.opendocument import load as odf_load
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation

from mcp_agent_rag.rag.archive_extractor import ArchiveExtractor
from mcp_agent_rag.utils import get_logger

logger = get_logger(__name__)

# Global OCR reader instance (lazy loaded)
_ocr_reader = None
# Sentinel object to indicate OCR initialization failure
_OCR_UNAVAILABLE = object()


class DocumentExtractor:
    """Extract text from various document formats."""

    SUPPORTED_EXTENSIONS = {
        ".txt", ".md", ".py", ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs",
        ".java", ".js", ".ts", ".sh", ".bat", ".ps1", ".s", ".asm",
        ".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp", ".pdf", ".html", ".htm",
        ".zip", ".7z", ".gz", ".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tbz2",
        ".tar.xz", ".txz", ".rar",
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"
    }

    @staticmethod
    def is_supported(file_path: Path) -> bool:
        """Check if file format is supported.

        Args:
            file_path: Path to file

        Returns:
            True if supported
        """
        path_str = str(file_path).lower()
        # Check for double extensions like .tar.gz
        for ext in DocumentExtractor.SUPPORTED_EXTENSIONS:
            if path_str.endswith(ext):
                return True
        return False

    @staticmethod
    def _get_ocr_reader():
        """Get or initialize the OCR reader (lazy loading).
        
        Returns:
            EasyOCR reader instance or None if initialization fails
        """
        global _ocr_reader
        if _ocr_reader is None:
            try:
                import easyocr
                
                # Check GPU availability and configuration
                gpu_available = False
                gpu_enabled = True  # Default to enabled
                
                # Try to load config to check GPU settings
                try:
                    from mcp_agent_rag.config import Config
                    config = Config()
                    gpu_enabled = config.get("gpu_enabled", True)
                except Exception:
                    # If config not available, just use auto-detection
                    pass
                
                # Only use GPU if both available and enabled in config
                if gpu_enabled:
                    try:
                        import torch
                        gpu_available = torch.cuda.is_available()
                    except ImportError:
                        pass
                
                logger.info("Initializing EasyOCR reader (this may take a moment)...")
                
                if gpu_available and gpu_enabled:
                    logger.info("GPU detected and enabled, using GPU acceleration for OCR")
                    _ocr_reader = easyocr.Reader(['en'], gpu=True, verbose=False)
                else:
                    if not gpu_enabled:
                        logger.info("GPU disabled in config, using CPU for OCR")
                    else:
                        logger.info("No GPU detected, using CPU for OCR")
                    # Suppress the pin_memory warning when using CPU-only mode
                    # This warning occurs because EasyOCR's internal DataLoader uses pin_memory=True
                    # by default, but pin_memory is only beneficial when using GPU acceleration.
                    # Since we're using gpu=False, the warning is harmless and can be safely ignored.
                    with warnings.catch_warnings():
                        warnings.filterwarnings(
                            'ignore',
                            message=".*pin_memory.*no accelerator.*",
                            category=UserWarning
                        )
                        _ocr_reader = easyocr.Reader(['en'], gpu=False, verbose=False)
                
                logger.info("EasyOCR reader initialized successfully")
            except Exception as e:
                logger.error(f"Failed to initialize EasyOCR: {e}")
                logger.error("Image text extraction will be unavailable")
                _ocr_reader = _OCR_UNAVAILABLE  # Mark as failed to avoid retrying
        return _ocr_reader if _ocr_reader is not _OCR_UNAVAILABLE else None

    @staticmethod
    def extract_text(file_path: Path) -> str | None:
        """Extract text from file.

        Args:
            file_path: Path to file

        Returns:
            Extracted text or None if failed
        """
        try:
            suffix = file_path.suffix.lower()

            if suffix in {".txt", ".md", ".py", ".c", ".cpp", ".h", ".hpp", ".cs",
                          ".go", ".rs", ".java", ".js", ".ts", ".sh", ".bat", ".ps1",
                          ".s", ".asm", ".cmake"}:
                return DocumentExtractor._extract_text_file(file_path)
            elif suffix == ".docx":
                return DocumentExtractor._extract_docx(file_path)
            elif suffix == ".xlsx":
                return DocumentExtractor._extract_xlsx(file_path)
            elif suffix == ".pptx":
                return DocumentExtractor._extract_pptx(file_path)
            elif suffix == ".odt":
                return DocumentExtractor._extract_odt(file_path)
            elif suffix == ".ods":
                return DocumentExtractor._extract_ods(file_path)
            elif suffix == ".odp":
                return DocumentExtractor._extract_odp(file_path)
            elif suffix == ".pdf":
                return DocumentExtractor._extract_pdf(file_path)
            elif suffix in {".html", ".htm"}:
                return DocumentExtractor._extract_html(file_path)
            elif suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"}:
                return DocumentExtractor._extract_image(file_path)
            else:
                logger.warning(f"Unsupported file format: {suffix}")
                return None
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return None

    @staticmethod
    def _extract_text_file(file_path: Path) -> str:
        """Extract text from plain text file."""
        try:
            with open(file_path, "rb") as f:
                raw_data = f.read()
                result = chardet.detect(raw_data)
                encoding = result["encoding"] or "utf-8"
            with open(file_path, encoding=encoding, errors="replace") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return None

    @staticmethod
    def _extract_docx(file_path: Path) -> str:
        """Extract text from DOCX file, including OCR on embedded images."""
        doc = DocxDocument(str(file_path))
        text_parts = []
        
        # Extract regular text
        for para in doc.paragraphs:
            text_parts.append(para.text)
        
        # Extract text from images if OCR is available
        if DocumentExtractor._get_ocr_reader() is not None:
            try:
                # Access document relationships to get images
                for rel_id, rel in doc.part.rels.items():
                    if "image" in rel.target_ref:
                        try:
                            image_data = rel.target_part.blob
                            image_text = DocumentExtractor._extract_text_from_image_bytes(
                                image_data,
                                f"DOCX embedded image {rel_id}"
                            )
                            if image_text.strip():
                                text_parts.append(f"\n[Embedded Image]:\n{image_text}")
                        except Exception as e:
                            logger.debug(f"Could not extract text from embedded image: {e}")
            except Exception as e:
                logger.debug(f"Could not access images in DOCX: {e}")
        
        return "\n".join(text_parts)

    @staticmethod
    def _extract_xlsx(file_path: Path) -> str:
        """Extract text from XLSX file."""
        wb = load_workbook(str(file_path), read_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}\n")
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts)

    @staticmethod
    def _extract_pptx(file_path: Path) -> str:
        """Extract text from PPTX file, including OCR on embedded images."""
        prs = Presentation(str(file_path))
        text_parts = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"Slide {i + 1}:")
            for shape in slide.shapes:
                # Extract text from shapes
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
                
                # Extract text from images if OCR is available
                if DocumentExtractor._get_ocr_reader() is not None:
                    if hasattr(shape, "image"):
                        try:
                            image_data = shape.image.blob
                            image_text = DocumentExtractor._extract_text_from_image_bytes(
                                image_data,
                                f"PPTX slide {i + 1} image"
                            )
                            if image_text.strip():
                                text_parts.append(f"[Image]:\n{image_text}")
                        except Exception as e:
                            logger.debug(f"Could not extract text from slide image: {e}")
        return "\n".join(text_parts)

    @staticmethod
    def _extract_odt(file_path: Path) -> str:
        """Extract text from ODT file."""
        doc = odf_load(str(file_path))
        paragraphs = doc.getElementsByType(odf_text.P)
        return "\n".join([teletype.extractText(p) for p in paragraphs])

    @staticmethod
    def _extract_ods(file_path: Path) -> str:
        """Extract text from ODS file."""
        doc = odf_load(str(file_path))
        text_parts = []
        from odf.table import Table, TableCell, TableRow
        tables = doc.getElementsByType(Table)
        for table in tables:
            text_parts.append(f"Table: {table.getAttribute('name') or 'Unnamed'}")
            rows = table.getElementsByType(TableRow)
            for row in rows:
                cells = row.getElementsByType(TableCell)
                row_text = "\t".join([teletype.extractText(cell) for cell in cells])
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts)

    @staticmethod
    def _extract_odp(file_path: Path) -> str:
        """Extract text from ODP file."""
        doc = odf_load(str(file_path))
        from odf.draw import Page as DrawPage
        pages = doc.getElementsByType(DrawPage)
        text_parts = []
        for i, page in enumerate(pages):
            text_parts.append(f"Slide {i + 1}:")
            paragraphs = page.getElementsByType(odf_text.P)
            for p in paragraphs:
                text = teletype.extractText(p)
                if text.strip():
                    text_parts.append(text)
        return "\n".join(text_parts)

    @staticmethod
    def _extract_pdf(file_path: Path) -> str:
        """Extract text from PDF file, including OCR on embedded images."""
        import logging
        
        # Suppress pypdf warnings about malformed PDFs
        # These warnings (Invalid Lookup Table, Fax4Decode errors, image mask size mismatches)
        # are non-fatal and occur when processing PDFs with corrupted compression streams,
        # invalid lookup tables, or misaligned image masks. pypdf continues processing despite
        # these issues, so we suppress the warnings to avoid cluttering the logs.
        pypdf_logger = logging.getLogger("pypdf")
        original_level = pypdf_logger.level
        pypdf_logger.setLevel(logging.ERROR)
        
        try:
            with open(file_path, "rb") as f:
                # Suppress warnings about PDF parsing issues
                with warnings.catch_warnings():
                    warnings.filterwarnings(
                        'ignore',
                        message=".*Invalid Lookup Table.*",
                        category=Warning
                    )
                    warnings.filterwarnings(
                        'ignore',
                        message=".*Fax4Decode.*",
                        category=Warning
                    )
                    warnings.filterwarnings(
                        'ignore',
                        message=".*image and mask.*",
                        category=Warning
                    )
                    
                    reader = pypdf.PdfReader(f)
                    text_parts = []
                    for i, page in enumerate(reader.pages):
                        # Extract regular text
                        text = page.extract_text()
                        if text.strip():
                            text_parts.append(f"Page {i + 1}:\n{text}")
                        
                        # Extract images and perform OCR if available
                        if hasattr(page, 'images') and DocumentExtractor._get_ocr_reader() is not None:
                            for img_index, image in enumerate(page.images):
                                try:
                                    image_text = DocumentExtractor._extract_text_from_image_bytes(
                                        image.data,
                                        f"PDF page {i + 1} image {img_index + 1}"
                                    )
                                    if image_text.strip():
                                        text_parts.append(
                                            f"Page {i + 1} - Image {img_index + 1}:\n{image_text}"
                                        )
                                except Exception as e:
                                    logger.debug(f"Could not extract text from image on page {i + 1}: {e}")
                                
                    return "\n".join(text_parts)
        finally:
            # Restore original logging level
            pypdf_logger.setLevel(original_level)

    @staticmethod
    def _extract_html(file_path: Path) -> str:
        """Extract text from HTML file."""
        with open(file_path, encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            return soup.get_text(separator="\n", strip=True)

    @staticmethod
    def _extract_image(file_path: Path) -> str:
        """Extract text from image file using OCR.
        
        Args:
            file_path: Path to image file
            
        Returns:
            Extracted text from image or empty string if no text found
        """
        reader = DocumentExtractor._get_ocr_reader()
        if reader is None:
            logger.warning(f"OCR not available, cannot extract text from {file_path}")
            return ""
        
        try:
            # Read image
            img = Image.open(file_path)
            
            # Convert to RGB if necessary (EasyOCR works best with RGB)
            if img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')
            
            # Perform OCR
            results = reader.readtext(str(file_path), detail=0)
            
            # Combine results
            if results:
                extracted_text = "\n".join(results)
                logger.debug(f"Extracted {len(results)} text blocks from {file_path.name}")
                return extracted_text
            else:
                logger.debug(f"No text found in image {file_path.name}")
                return ""
                
        except Exception as e:
            logger.error(f"Error extracting text from image {file_path}: {e}")
            return ""

    @staticmethod
    def _extract_text_from_image_bytes(image_bytes: bytes, source_name: str = "embedded") -> str:
        """Extract text from image bytes using OCR.
        
        Args:
            image_bytes: Image data as bytes
            source_name: Name/description of the image source for logging
            
        Returns:
            Extracted text from image or empty string if no text found
        """
        reader = DocumentExtractor._get_ocr_reader()
        if reader is None:
            return ""
        
        try:
            # Open image from bytes
            img = Image.open(io.BytesIO(image_bytes))
            
            # Convert to RGB if necessary
            if img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')
            
            # Convert to numpy array for EasyOCR
            img_array = np.array(img)
            
            # Perform OCR
            results = reader.readtext(img_array, detail=0)
            
            # Combine results
            if results:
                extracted_text = "\n".join(results)
                logger.debug(f"Extracted {len(results)} text blocks from {source_name}")
                return extracted_text
            else:
                logger.debug(f"No text found in {source_name}")
                return ""
                
        except Exception as e:
            logger.error(f"Error extracting text from {source_name}: {e}")
            return ""


def find_files_to_process(
    path: str,
    recursive: bool = False,
    glob_pattern: str | None = None,
    respect_gitignore: bool = True,
) -> list[Path]:
    """Find files to process, including extracting archives.

    Args:
        path: Path to file or directory
        recursive: Whether to recurse into subdirectories
        glob_pattern: Optional glob pattern to match files
        respect_gitignore: Whether to respect .gitignore

    Returns:
        List of file paths to process (including files extracted from archives)
    """
    path_obj = Path(path).expanduser().resolve()
    files = []

    # Load gitignore patterns if requested
    gitignore_patterns = []
    if respect_gitignore:
        parent_dir = path_obj if path_obj.is_dir() else path_obj.parent
        gitignore_patterns = _load_gitignore_patterns(parent_dir)

    if path_obj.is_file():
        if ArchiveExtractor.is_archive(path_obj):
            # Extract archive and get files
            logger.info(f"Extracting archive: {path_obj.name}")
            extracted_files = ArchiveExtractor.extract_archive(path_obj)
            # Filter to only supported document types
            for extracted_file in extracted_files:
                is_supported = DocumentExtractor.is_supported(extracted_file)
                is_not_archive = not ArchiveExtractor.is_archive(extracted_file)
                if is_supported and is_not_archive:
                    files.append(extracted_file)
        elif DocumentExtractor.is_supported(path_obj):
            files.append(path_obj)
    elif path_obj.is_dir():
        if glob_pattern:
            # Use glob pattern
            if recursive:
                candidate_files = path_obj.rglob(glob_pattern)
            else:
                candidate_files = path_obj.glob(glob_pattern)

            for file_path in candidate_files:
                if file_path.is_file():
                    if ArchiveExtractor.is_archive(file_path):
                        # Extract archive and get files
                        logger.info(f"Extracting archive: {file_path.name}")
                        extracted_files = ArchiveExtractor.extract_archive(file_path)
                        # Filter to only supported document types
                        for extracted_file in extracted_files:
                            is_supported = DocumentExtractor.is_supported(extracted_file)
                            is_not_archive = not ArchiveExtractor.is_archive(
                                extracted_file
                            )
                            if is_supported and is_not_archive:
                                files.append(extracted_file)
                    else:
                        files.append(file_path)
        else:
            # Get all supported files
            if recursive:
                for root, dirs, filenames in os.walk(path_obj):
                    root_path = Path(root)
                    # Filter directories based on gitignore
                    if gitignore_patterns:
                        filtered_dirs = [
                            d for d in dirs
                            if not _is_ignored(root_path / d, gitignore_patterns)
                        ]
                        dirs[:] = filtered_dirs
                    for filename in filenames:
                        file_path = root_path / filename
                        if gitignore_patterns and _is_ignored(
                            file_path, gitignore_patterns
                        ):
                            continue

                        if ArchiveExtractor.is_archive(file_path):
                            # Extract archive and get files
                            logger.info(f"Extracting archive: {file_path.name}")
                            extracted_files = ArchiveExtractor.extract_archive(
                                file_path
                            )
                            # Filter to only supported document types
                            for extracted_file in extracted_files:
                                is_supported = DocumentExtractor.is_supported(
                                    extracted_file
                                )
                                is_not_archive = not ArchiveExtractor.is_archive(
                                    extracted_file
                                )
                                if is_supported and is_not_archive:
                                    files.append(extracted_file)
                        elif DocumentExtractor.is_supported(file_path):
                            files.append(file_path)
            else:
                for file_path in path_obj.iterdir():
                    if not file_path.is_file():
                        continue
                    if gitignore_patterns and _is_ignored(file_path, gitignore_patterns):
                        continue

                    if ArchiveExtractor.is_archive(file_path):
                        # Extract archive and get files
                        logger.info(f"Extracting archive: {file_path.name}")
                        extracted_files = ArchiveExtractor.extract_archive(file_path)
                        # Filter to only supported document types
                        for extracted_file in extracted_files:
                            is_supported = DocumentExtractor.is_supported(
                                extracted_file
                            )
                            is_not_archive = not ArchiveExtractor.is_archive(
                                extracted_file
                            )
                            if is_supported and is_not_archive:
                                files.append(extracted_file)
                    elif DocumentExtractor.is_supported(file_path):
                        files.append(file_path)

    # Filter out ignored files
    if gitignore_patterns:
        files = [f for f in files if not _is_ignored(f, gitignore_patterns)]

    return sorted(files)


def _load_gitignore_patterns(directory: Path) -> list[str]:
    """Load .gitignore patterns from directory.

    Args:
        directory: Directory to search for .gitignore

    Returns:
        List of patterns
    """
    patterns = []
    gitignore_path = directory / ".gitignore"
    if gitignore_path.exists():
        with open(gitignore_path, encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if line and not line.startswith("#"):
                    patterns.append(line)
    return patterns


def _is_ignored(path: Path, patterns: list[str]) -> bool:
    """Check if path matches any gitignore pattern.

    Args:
        path: Path to check
        patterns: List of gitignore patterns

    Returns:
        True if path should be ignored
    """
    path_str = str(path)
    for pattern in patterns:
        # Simple pattern matching (simplified version)
        if pattern.startswith("/"):
            pattern = pattern[1:]
        if pattern.endswith("/"):
            if path.is_dir() and path.name == pattern[:-1]:
                return True
        elif "*" in pattern:
            import fnmatch
            if fnmatch.fnmatch(path.name, pattern):
                return True
        elif path.name == pattern:
            return True
        elif pattern in path_str:
            return True
    return False
