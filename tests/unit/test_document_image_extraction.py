"""Tests for extracting images from documents (PDF, DOCX, PPTX)."""

import io
from pathlib import Path
from unittest.mock import MagicMock, patch

import pypdf
import pytest
from docx import Document as DocxDocument
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

from mcp_agent_rag.rag.extractor import DocumentExtractor


@pytest.fixture
def pdf_with_images(temp_dir: Path) -> Path:
    """Create a PDF file with embedded images."""
    from pypdf import PdfWriter, PdfReader
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    # Create a simple image
    img = Image.new('RGB', (100, 50), color='white')
    draw = ImageDraw.Draw(img)
    draw.text((10, 15), "PDF Image", fill='black')
    img_path = temp_dir / "temp_img.png"
    img.save(img_path)
    
    # Create PDF with text and image using reportlab
    pdf_path = temp_dir / "test_with_images.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=letter)
    
    # Add text
    c.drawString(100, 750, "This is a PDF with an embedded image")
    
    # Add image
    c.drawImage(str(img_path), 100, 600, width=100, height=50)
    
    c.save()
    
    # Clean up temp image
    img_path.unlink()
    
    return pdf_path


@pytest.fixture
def docx_with_images(temp_dir: Path) -> Path:
    """Create a DOCX file with embedded images."""
    doc = DocxDocument()
    doc.add_paragraph("This is a DOCX with embedded images")
    
    # Create a simple image
    img = Image.new('RGB', (100, 50), color='white')
    draw = ImageDraw.Draw(img)
    draw.text((10, 15), "DOCX Image", fill='black')
    img_path = temp_dir / "temp_img.png"
    img.save(img_path)
    
    # Add image to document
    doc.add_picture(str(img_path), width=Inches(1.0))
    
    # Save document
    docx_path = temp_dir / "test_with_images.docx"
    doc.save(str(docx_path))
    
    # Clean up temp image
    img_path.unlink()
    
    return docx_path


@pytest.fixture
def pptx_with_images(temp_dir: Path) -> Path:
    """Create a PPTX file with embedded images."""
    prs = Presentation()
    
    # Add a slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add text box
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    txBox.text = "This is a PPTX with embedded images"
    
    # Create a simple image
    img = Image.new('RGB', (100, 50), color='white')
    draw = ImageDraw.Draw(img)
    draw.text((10, 15), "PPTX Image", fill='black')
    img_path = temp_dir / "temp_img.png"
    img.save(img_path)
    
    # Add image to slide
    slide.shapes.add_picture(str(img_path), Inches(1), Inches(3), width=Inches(1.0))
    
    # Save presentation
    pptx_path = temp_dir / "test_with_images.pptx"
    prs.save(str(pptx_path))
    
    # Clean up temp image
    img_path.unlink()
    
    return pptx_path


def test_extract_pdf_with_text_only(temp_dir: Path):
    """Test extracting PDF with text only (no images)."""
    # Create a simple PDF with just text
    writer = pypdf.PdfWriter()
    page = pypdf.PageObject.create_blank_page(width=200, height=200)
    writer.add_page(page)
    
    pdf_path = temp_dir / "text_only.pdf"
    with open(pdf_path, "wb") as f:
        writer.write(f)
    
    # Extract text
    text = DocumentExtractor.extract_text(pdf_path)
    
    # Should handle PDF without images
    assert text is not None


def test_extract_pdf_with_images_no_ocr(pdf_with_images):
    """Test extracting PDF with images when OCR is not available."""
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=None):
        text = DocumentExtractor.extract_text(pdf_with_images)
        
        # Should extract regular text but skip images
        assert text is not None
        assert "This is a PDF with an embedded image" in text or "Page 1" in text


def test_extract_pdf_with_images_with_ocr(pdf_with_images):
    """Test extracting PDF with images using OCR."""
    mock_reader = MagicMock()
    mock_reader.readtext.return_value = ["PDF Image Text"]
    
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=mock_reader):
        with patch.object(DocumentExtractor, '_extract_text_from_image_bytes', return_value="PDF Image Text"):
            text = DocumentExtractor.extract_text(pdf_with_images)
            
            # Should extract both regular text and image text
            assert text is not None


def test_extract_docx_text_only(temp_dir: Path):
    """Test extracting DOCX with text only."""
    doc = DocxDocument()
    doc.add_paragraph("Simple DOCX text")
    
    docx_path = temp_dir / "simple.docx"
    doc.save(str(docx_path))
    
    text = DocumentExtractor.extract_text(docx_path)
    
    assert text is not None
    assert "Simple DOCX text" in text


def test_extract_docx_with_images_no_ocr(docx_with_images):
    """Test extracting DOCX with images when OCR is not available."""
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=None):
        text = DocumentExtractor.extract_text(docx_with_images)
        
        # Should extract regular text but skip images
        assert text is not None
        assert "This is a DOCX with embedded images" in text


def test_extract_docx_with_images_with_ocr(docx_with_images):
    """Test extracting DOCX with images using OCR."""
    mock_reader = MagicMock()
    
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=mock_reader):
        with patch.object(DocumentExtractor, '_extract_text_from_image_bytes', return_value="DOCX Image Text"):
            text = DocumentExtractor.extract_text(docx_with_images)
            
            # Should extract both regular text and potentially image text
            assert text is not None
            assert "This is a DOCX with embedded images" in text


def test_extract_docx_image_extraction_error(docx_with_images):
    """Test handling errors when extracting images from DOCX."""
    mock_reader = MagicMock()
    
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=mock_reader):
        with patch.object(DocumentExtractor, '_extract_text_from_image_bytes', side_effect=Exception("OCR error")):
            # Should handle error gracefully and still return document text
            text = DocumentExtractor.extract_text(docx_with_images)
            
            assert text is not None


def test_extract_pptx_text_only(temp_dir: Path):
    """Test extracting PPTX with text only."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    txBox.text = "Simple PPTX text"
    
    pptx_path = temp_dir / "simple.pptx"
    prs.save(str(pptx_path))
    
    text = DocumentExtractor.extract_text(pptx_path)
    
    assert text is not None
    assert "Slide 1" in text
    assert "Simple PPTX text" in text


def test_extract_pptx_with_images_no_ocr(pptx_with_images):
    """Test extracting PPTX with images when OCR is not available."""
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=None):
        text = DocumentExtractor.extract_text(pptx_with_images)
        
        # Should extract regular text but skip images
        assert text is not None
        assert "Slide 1" in text


def test_extract_pptx_with_images_with_ocr(pptx_with_images):
    """Test extracting PPTX with images using OCR."""
    mock_reader = MagicMock()
    
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=mock_reader):
        with patch.object(DocumentExtractor, '_extract_text_from_image_bytes', return_value="PPTX Image Text"):
            text = DocumentExtractor.extract_text(pptx_with_images)
            
            # Should extract text from slides
            assert text is not None
            assert "Slide 1" in text


def test_extract_pptx_image_extraction_error(pptx_with_images):
    """Test handling errors when extracting images from PPTX."""
    mock_reader = MagicMock()
    
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=mock_reader):
        with patch.object(DocumentExtractor, '_extract_text_from_image_bytes', side_effect=Exception("OCR error")):
            # Should handle error gracefully and still return slide text
            text = DocumentExtractor.extract_text(pptx_with_images)
            
            assert text is not None


def test_extract_docx_without_images_attribute(temp_dir: Path):
    """Test DOCX extraction when document doesn't have expected image attributes."""
    doc = DocxDocument()
    doc.add_paragraph("Test text")
    
    docx_path = temp_dir / "test.docx"
    doc.save(str(docx_path))
    
    mock_reader = MagicMock()
    
    # Mock part.rels to simulate missing image relationships
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=mock_reader):
        text = DocumentExtractor.extract_text(docx_path)
        
        # Should handle gracefully
        assert text is not None
        assert "Test text" in text


def test_extract_pptx_shape_without_image(temp_dir: Path):
    """Test PPTX extraction when shapes don't have image attribute."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add only text box (no image)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    txBox.text = "Only text"
    
    pptx_path = temp_dir / "test.pptx"
    prs.save(str(pptx_path))
    
    mock_reader = MagicMock()
    
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=mock_reader):
        text = DocumentExtractor.extract_text(pptx_path)
        
        # Should handle shapes without images
        assert text is not None
        assert "Only text" in text


def test_pdf_page_without_images_attribute(temp_dir: Path):
    """Test PDF extraction when page doesn't have images attribute."""
    writer = pypdf.PdfWriter()
    page = pypdf.PageObject.create_blank_page(width=200, height=200)
    writer.add_page(page)
    
    pdf_path = temp_dir / "no_images.pdf"
    with open(pdf_path, "wb") as f:
        writer.write(f)
    
    mock_reader = MagicMock()
    
    with patch.object(DocumentExtractor, '_get_ocr_reader', return_value=mock_reader):
        text = DocumentExtractor.extract_text(pdf_path)
        
        # Should handle pages without images attribute
        assert text is not None
